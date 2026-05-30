"""
批量文字替换工具 - 对文件夹内所有可编辑文件进行文字替换。
先复制源文件夹，在副本上操作，保留原文件格式。
支持: Word(.doc/.docx), Excel(.xls/.xlsx), PPT(.pptx), 纯文本文件
"""

import os
import shutil
import tempfile
import threading
import tkinter as tk
from datetime import datetime
from tkinter import filedialog, messagebox, ttk
from pathlib import Path

# ── 可处理的文本文件扩展名 ──────────────────────────────────────────
TEXT_EXTENSIONS = {
    ".txt", ".csv", ".json", ".xml", ".html", ".htm", ".md", ".markdown",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".css", ".scss", ".less",
    ".java", ".c", ".cpp", ".h", ".hpp", ".cs", ".go", ".rs", ".rb",
    ".php", ".swift", ".kt", ".scala", ".r", ".m", ".sh", ".bat",
    ".ps1", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".sql", ".log", ".tex", ".rst", ".svg",
}

SKIP_EXTENSIONS = {
    ".pdf", ".exe", ".dll", ".so", ".dylib",
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".ico", ".tiff", ".webp",
    ".mp3", ".mp4", ".avi", ".mkv", ".mov", ".wmv", ".flv", ".wav",
    ".zip", ".rar", ".7z", ".tar", ".gz", ".bz2",
    ".ppt",  # 旧格式 ppt 暂不支持
    ".pyc", ".class", ".o", ".obj", ".lib",
    ".ttf", ".otf", ".woff", ".woff2", ".eot",
    ".db", ".sqlite", ".sqlite3", ".mdb",
    ".bin", ".dat", ".pkl", ".pickle",
}


# ── 核心替换逻辑 ────────────────────────────────────────────────────

def copy_folder(src, dst):
    """复制源文件夹到目标路径，dst 必须不存在。"""
    shutil.copytree(src, dst)


def replace_in_text_file(filepath, find_text, replace_text):
    """对纯文本文件执行替换，返回替换次数。"""
    try:
        content = Path(filepath).read_text(encoding="utf-8")
    except (UnicodeDecodeError, UnicodeError):
        try:
            content = Path(filepath).read_text(encoding="gbk")
        except Exception:
            return 0

    if find_text not in content:
        return 0

    count = content.count(find_text)
    new_content = content.replace(find_text, replace_text)

    enc = "utf-8"
    try:
        content.encode("gbk")
        new_content.encode("gbk")
        enc = "gbk"
    except Exception:
        pass

    Path(filepath).write_text(new_content, encoding=enc)
    return count


def _date_to_iso(value):
    """日期类型（datetime / COM TimeType 等）→ ISO 字符串，失败返回 None。"""
    try:
        return (f"{value.year:04d}-{value.month:02d}-{value.day:02d} "
                f"{getattr(value, 'hour', 0):02d}:"
                f"{getattr(value, 'minute', 0):02d}:"
                f"{getattr(value, 'second', 0):02d}")
    except (AttributeError, TypeError):
        return None


def _serial_to_datetime(serial):
    """Excel 日期序列号 → datetime（1900 日期系统）。"""
    from datetime import timedelta
    if serial <= 0:
        return None
    base = datetime(1899, 12, 30)
    try:
        return base + timedelta(days=float(serial))
    except (ValueError, OverflowError):
        return None


def _iso_to_datetime(iso_str):
    """ISO 字符串 → datetime，失败返回 None。"""
    try:
        parts = iso_str.strip().split(' ')
        d = parts[0].split('-')
        t = parts[1].split(':') if len(parts) > 1 else ['0', '0', '0']
        return datetime(
            int(d[0]), int(d[1]), int(d[2]),
            int(t[0]), int(t[1]), int(t[2]),
        )
    except (ValueError, IndexError, TypeError, AttributeError):
        return None


def _replace_in_runs(runs, find_text, replace_text):
    """替换一段 runs 中的文字，处理中文跨 run 拆分的情况。返回替换次数。"""
    full_text = "".join(run.text or "" for run in runs)
    if find_text not in full_text:
        return 0

    occurrence_count = full_text.count(find_text)

    # 先尝试逐 run 替换（可保留格式）
    for run in runs:
        if find_text in (run.text or ""):
            run.text = run.text.replace(find_text, replace_text)

    # 检查是否全部替换成功
    remaining = "".join(run.text or "" for run in runs)
    if find_text not in remaining:
        return occurrence_count

    # 有关键字跨 run 拆分，段落级替换：全量文字放入第一个 run
    new_text = full_text.replace(find_text, replace_text)
    if runs:
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ""
    return occurrence_count


def replace_in_docx(filepath, find_text, replace_text):
    """对 .docx 文件执行替换，返回替换次数。
    先通过 python-docx 高层 API 处理正文/表格/页眉页脚，
    再通过 XML 级别扫描覆盖内容控件(SDT)、文本框等被遗漏的结构。"""
    try:
        from docx import Document
    except ImportError:
        return -1

    doc = Document(filepath)
    count = 0

    for para in doc.paragraphs:
        count += _replace_in_runs(para.runs, find_text, replace_text)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    count += _replace_in_runs(para.runs, find_text, replace_text)
                # 嵌套表格
                for nested_table in cell.tables:
                    for nrow in nested_table.rows:
                        for ncell in nrow.cells:
                            for para in ncell.paragraphs:
                                count += _replace_in_runs(para.runs, find_text, replace_text)

    # 页眉页脚（仅处理已存在的，避免 python-docx 自动创建空页眉页脚）
    NS = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'

    def _has_header(section):
        for child in section._sectPr:
            if child.tag == f'{{{NS}}}headerReference':
                return True
        return False

    def _has_footer(section):
        for child in section._sectPr:
            if child.tag == f'{{{NS}}}footerReference':
                return True
        return False

    for section in doc.sections:
        if _has_header(section):
            for para in section.header.paragraphs:
                count += _replace_in_runs(para.runs, find_text, replace_text)
        if _has_footer(section):
            for para in section.footer.paragraphs:
                count += _replace_in_runs(para.runs, find_text, replace_text)

    # XML 级别扫描：覆盖内容控件(SDT)、文本框等 python-docx 高层 API 遗漏的结构
    # 注：正文只扫 SDT 内部，避免破坏已处理段落的段落格式导致页数变化

    def _xml_scan_sdt(root_elem):
        """在 SDT（内容控件）内逐 w:t 替换。"""
        n = 0
        for sdt in root_elem.iter(f'{{{NS}}}sdt'):
            for elem in sdt.iter():
                if elem.tag == f'{{{NS}}}t' and elem.text and find_text in elem.text:
                    n += elem.text.count(find_text)
                    elem.text = elem.text.replace(find_text, replace_text)
        return n

    def _xml_scan_sdt_cross_run(root_elem):
        """在 SDT 段落内处理跨 w:t 拆分。"""
        n = 0
        for sdt in root_elem.iter(f'{{{NS}}}sdt'):
            for para_elem in sdt.iter(f'{{{NS}}}p'):
                wt_elems = para_elem.findall(f'.//{{{NS}}}t')
                if len(wt_elems) < 2:
                    continue
                combined = ''.join(wt.text or '' for wt in wt_elems)
                if find_text in combined:
                    n += combined.count(find_text)
                    wt_elems[0].text = combined.replace(find_text, replace_text)
                    for wt in wt_elems[1:]:
                        wt.text = ''
        return n

    count += _xml_scan_sdt(doc.element.body)
    count += _xml_scan_sdt_cross_run(doc.element.body)

    # 页眉页脚：完整 XML 扫描（结构小，不影响正文页数）
    for section in doc.sections:
        parts = []
        if _has_header(section):
            parts.append(section.header)
        if _has_footer(section):
            parts.append(section.footer)
        for part in parts:
            for elem in part._element.iter():
                if elem.tag == f'{{{NS}}}t' and elem.text and find_text in elem.text:
                    n = elem.text.count(find_text)
                    elem.text = elem.text.replace(find_text, replace_text)
                    count += n
            for para_elem in part._element.iter(f'{{{NS}}}p'):
                wt_elems = para_elem.findall(f'.//{{{NS}}}t')
                if len(wt_elems) >= 2:
                    combined = ''.join(wt.text or '' for wt in wt_elems)
                    if find_text in combined:
                        n = combined.count(find_text)
                        wt_elems[0].text = combined.replace(find_text, replace_text)
                        for wt in wt_elems[1:]:
                            wt.text = ''
                        count += n

    if count:
        doc.save(filepath)
    return count


def replace_in_xlsx(filepath, find_text, replace_text):
    """对 Excel 文件执行替换，只处理字符串单元格，忽略公式，返回替换次数。"""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return -1

    wb = load_workbook(filepath)
    count = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is None:
                    continue
                # 跳过公式单元格
                if isinstance(cell.value, str) and cell.value.startswith("="):
                    continue
                if isinstance(cell.value, str) and find_text in cell.value:
                    cell.value = cell.value.replace(find_text, replace_text)
                    count += 1
                elif isinstance(cell.value, (int, float)):
                    s = str(int(cell.value)) if isinstance(cell.value, float) and cell.value == int(cell.value) else str(cell.value)
                    if find_text in s:
                        new_s = s.replace(find_text, replace_text)
                        try:
                            cell.value = int(new_s) if '.' not in new_s else float(new_s)
                            count += 1
                        except (ValueError, TypeError):
                            pass
                else:
                    iso = _date_to_iso(cell.value)
                    if iso and find_text in iso:
                        new_val = _iso_to_datetime(iso.replace(find_text, replace_text))
                        if new_val is not None:
                            cell.value = new_val
                            count += 1

    if count:
        wb.save(filepath)
    else:
        wb.close()
    return count


def replace_in_pptx(filepath, find_text, replace_text):
    """对 .pptx 文件执行替换，返回替换次数。"""
    try:
        from pptx import Presentation
    except ImportError:
        return -1

    prs = Presentation(filepath)
    count = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    count += _replace_in_runs(para.runs, find_text, replace_text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            count += _replace_in_runs(para.runs, find_text, replace_text)

    if count:
        prs.save(filepath)
    return count


# ── COM 单例（避免反复创建/销毁导致挂起）──────────────────────────

_word = None
_excel = None


def _get_word():
    global _word
    if _word is None:
        import win32com.client
        _word = win32com.client.Dispatch("Word.Application")
        _word.Visible = False
        _word.DisplayAlerts = 0
        _word.ScreenUpdating = False
    return _word


def _get_excel():
    global _excel
    if _excel is None:
        import win32com.client
        _excel = win32com.client.Dispatch("Excel.Application")
        _excel.Visible = False
        _excel.DisplayAlerts = False
        _excel.ScreenUpdating = False
        _excel.EnableEvents = False
    return _excel


def _cleanup_com():
    global _word, _excel
    for obj in (_word, _excel):
        if obj is not None:
            try:
                obj.Quit()
            except Exception:
                pass
    _word = None
    _excel = None


def replace_in_doc(filepath, find_text, replace_text):
    """对 .doc 旧格式文件执行替换（通过 Word COM），返回替换次数。
    遇到路径包含特殊字符导致 Word 无法打开时，自动复制到临时目录处理。"""
    try:
        import win32com.client
    except ImportError:
        return -1

    filepath = str(Path(filepath).resolve())

    def _open_and_replace(path):
        doc = word.Documents.Open(path)
        try:
            body_text = doc.Content.Text

            def _do_replace(rng):
                f = rng.Find
                f.ClearFormatting()
                f.Replacement.ClearFormatting()
                f.Execute(
                    find_text, False, False, False, False, False,
                    True, 1, False, replace_text, 2,
                )

            # 正文替换
            body_count = body_text.count(find_text) if find_text in body_text else 0
            if body_count > 0:
                word.Selection.HomeKey(Unit=6)  # wdStory
                _do_replace(word.Selection)

            # 页眉页脚替换
            for section in doc.Sections:
                for hf_type in (1, 2, 3):  # Primary, FirstPage, EvenPages
                    try:
                        header = section.Headers(hf_type)
                        if header.Exists and find_text in header.Range.Text:
                            _do_replace(header.Range)
                    except Exception:
                        pass
                    try:
                        footer = section.Footers(hf_type)
                        if footer.Exists and find_text in footer.Range.Text:
                            _do_replace(footer.Range)
                    except Exception:
                        pass

            doc.Save()

            # 计数：正文 + 页眉页脚
            if body_count > 0:
                return body_count
            for section in doc.Sections:
                for hf_type in (1, 2, 3):
                    try:
                        if section.Headers(hf_type).Exists:
                            t = section.Headers(hf_type).Range.Text
                            if find_text in t:
                                return t.count(find_text) or 1
                    except Exception:
                        pass
                    try:
                        if section.Footers(hf_type).Exists:
                            t = section.Footers(hf_type).Range.Text
                            if find_text in t:
                                return t.count(find_text) or 1
                    except Exception:
                        pass
            return 0
        finally:
            doc.Close()

    word = _get_word()

    # Try direct open first
    try:
        return _open_and_replace(filepath)
    except Exception:
        pass

    # Fallback: the path may contain chars Word COM can't handle
    # Copy to a temp dir with a simple ASCII name, process, copy back
    tmpdir = tempfile.mkdtemp()
    try:
        ext = os.path.splitext(filepath)[1]
        tmpfile = os.path.join(tmpdir, f"_tmp{ext}")
        shutil.copy2(filepath, tmpfile)
        count = _open_and_replace(tmpfile)
        if count:
            shutil.copy2(tmpfile, filepath)
        return count
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)


def replace_in_xls(filepath, find_text, replace_text):
    """对 .xls 旧格式文件执行替换，多层回退：Excel COM → xlrd+xlutils → xlrd+xlwt。"""

    # 调试日志
    def _dbg_write(msg):
        try:
            with open("D:/xls_debug.log", "a", encoding="utf-8") as _f:
                _f.write(f"[{os.path.basename(filepath)}] {msg}\n")
        except Exception:
            pass

    try:
        with open("D:/xls_debug.log", "a", encoding="utf-8") as _f:
            _f.write(f"=== {os.path.basename(filepath)} ===\n")
            _f.write(f"FIND: {find_text!r} -> {replace_text!r}\n")
    except Exception:
        pass

    # 方案一: Excel COM — 原生保留所有格式，最可靠
    try:
        import win32com.client
        excel = _get_excel()
        xls_path = str(Path(filepath).resolve())
        wb = excel.Workbooks.Open(xls_path)
        if wb is not None:
            count = 0
            try:
                for ws in wb.Worksheets:
                    used = ws.UsedRange
                    if used is None:
                        continue
                    for row in used.Rows:
                        for cell in row.Cells:
                            try:
                                val = cell.Value  # 只读一次，避免 COM 线程每次返回不同类型
                                if val is None:
                                    _dbg_write(f"CELL: val=None")
                                    continue
                                vt = type(val).__name__
                                _dbg_write(f"CELL: type={vt}, val={val!r:.100}")
                                if isinstance(val, str):
                                    if find_text in val:
                                        cell.Value = val.replace(find_text, replace_text)
                                        count += 1
                                        _dbg_write(f"CELL: str direct REPLACED")
                                    else:
                                        # 后台线程中日期序列号可能以字符串形式返回（如 "45722.0"）
                                        try:
                                            display = str(cell.Text)
                                            if find_text in display:
                                                serial = float(val)
                                                dt_val = _serial_to_datetime(serial)
                                                if dt_val:
                                                    iso = _date_to_iso(dt_val)
                                                    if iso and find_text in iso:
                                                        new_val = _iso_to_datetime(iso.replace(find_text, replace_text))
                                                        if new_val:
                                                            cell.Value = new_val
                                                            count += 1
                                                            _dbg_write(f"CELL: str serial REPLACED")
                                                        else:
                                                            _dbg_write(f"CELL: str serial new_val=None")
                                                    else:
                                                        _dbg_write(f"CELL: str serial find_text not in iso")
                                                else:
                                                    _dbg_write(f"CELL: str serial dt_val=None")
                                            else:
                                                _dbg_write(f"CELL: str display no match: {display!r:.100}")
                                        except Exception as e:
                                            _dbg_write(f"CELL: str fallback ERROR: {e}")
                                elif isinstance(val, (int, float)):
                                    s = str(int(val)) if isinstance(val, float) and val == int(val) else str(val)
                                    if find_text in s:
                                        new_s = s.replace(find_text, replace_text)
                                        try:
                                            cell.Value = int(new_s) if '.' not in new_s else float(new_s)
                                            count += 1
                                            _dbg_write(f"CELL: numeric direct REPLACED")
                                        except (ValueError, TypeError):
                                            pass
                                    else:
                                        # 数字本身不含查找文字 → 检查显示文本（日期序列号等）
                                        try:
                                            display = str(cell.Text)
                                            if find_text in display:
                                                dt_val = _serial_to_datetime(float(val))
                                                if dt_val:
                                                    iso = _date_to_iso(dt_val)
                                                    if iso and find_text in iso:
                                                        new_val = _iso_to_datetime(iso.replace(find_text, replace_text))
                                                        if new_val is not None:
                                                            tz = getattr(val, 'tzinfo', None)
                                                            if tz is not None:
                                                                new_val = new_val.replace(tzinfo=tz)
                                                            cell.Value = new_val
                                                            count += 1
                                                            _dbg_write(f"CELL: numeric serial REPLACED")
                                        except Exception as e:
                                            _dbg_write(f"CELL: numeric fallback ERROR: {e}")
                                else:
                                    iso = _date_to_iso(val)
                                    if iso and find_text in iso:
                                        new_val = _iso_to_datetime(iso.replace(find_text, replace_text))
                                        if new_val is not None:
                                            tz = getattr(val, 'tzinfo', None)
                                            if tz is not None:
                                                new_val = new_val.replace(tzinfo=tz)
                                            cell.Value = new_val
                                            count += 1
                                            _dbg_write(f"CELL: date REPLACED")
                                        else:
                                            _dbg_write(f"CELL: date new_val=None")
                                    else:
                                        _dbg_write(f"CELL: date no match: iso={iso!r}")
                            except Exception as e:
                                _dbg_write(f"CELL OUTER ERROR: {e}")
                if count:
                    wb.Save()
            finally:
                wb.Close()
            _dbg_write(f"COM OK: count={count}")
            return count
        else:
            _dbg_write("COM FAIL: wb is None")
    except ImportError:
        _dbg_write("COM FAIL: ImportError")
    except Exception as e:
        _dbg_write(f"COM FAIL: {e}")

    # 方案二: xlrd + xlutils（保留格式，可能丢失被替换单元格的格式）
    try:
        import xlrd
        from xlutils.copy import copy
        rb = xlrd.open_workbook(filepath, formatting_info=True)
        wb = copy(rb)
        count = 0
        for sheet_idx in range(rb.nsheets):
            rs = rb.sheet_by_index(sheet_idx)
            ws = wb.get_sheet(sheet_idx)
            for row in range(rs.nrows):
                for col in range(rs.ncols):
                    ctype = rs.cell_type(row, col)
                    value = rs.cell_value(row, col)
                    if ctype == xlrd.XL_CELL_TEXT:
                        if isinstance(value, str) and find_text in value:
                            ws.write(row, col, value.replace(find_text, replace_text))
                            count += 1
                    elif ctype == xlrd.XL_CELL_NUMBER:
                        s = str(int(value)) if value == int(value) else str(value)
                        if find_text in s:
                            new_s = s.replace(find_text, replace_text)
                            try:
                                ws.write(row, col, int(new_s) if '.' not in new_s else float(new_s))
                                count += 1
                            except (ValueError, TypeError):
                                pass
                    elif ctype == xlrd.XL_CELL_DATE:
                        try:
                            dt_tuple = xlrd.xldate_as_tuple(value, rb.datemode)
                            dt_val = datetime(*dt_tuple)
                            iso = _date_to_iso(dt_val)
                            if iso and find_text in iso:
                                new_dt = _iso_to_datetime(iso.replace(find_text, replace_text))
                                if new_dt:
                                    ws.write(row, col, new_dt)
                                    count += 1
                        except Exception:
                            pass
        if count:
            wb.save(filepath)
        _dbg_write(f"xlutils: count={count}")
        return count
    except ImportError:
        _dbg_write("xlutils FAIL: ImportError")
    except Exception as e:
        _dbg_write(f"xlutils FAIL: {e}")

    # 方案三: xlrd 只读 + xlwt 重写（丢失全部格式，最后手段）
    try:
        import xlrd
        import xlwt
        rb = xlrd.open_workbook(filepath, formatting_info=False)
        wb = xlwt.Workbook()
        count = 0
        for sheet_idx in range(rb.nsheets):
            rs = rb.sheet_by_index(sheet_idx)
            ws = wb.add_sheet(rs.name or f"Sheet{sheet_idx+1}")
            for row in range(rs.nrows):
                for col in range(rs.ncols):
                    ctype = rs.cell_type(row, col)
                    value = rs.cell_value(row, col)
                    if ctype == xlrd.XL_CELL_TEXT:
                        if isinstance(value, str) and find_text in value:
                            value = value.replace(find_text, replace_text)
                            count += 1
                    elif ctype == xlrd.XL_CELL_NUMBER:
                        s = str(int(value)) if value == int(value) else str(value)
                        if find_text in s:
                            new_s = s.replace(find_text, replace_text)
                            try:
                                value = int(new_s) if '.' not in new_s else float(new_s)
                                count += 1
                            except (ValueError, TypeError):
                                pass
                    elif ctype == xlrd.XL_CELL_DATE:
                        try:
                            dt_tuple = xlrd.xldate_as_tuple(value, rb.datemode)
                            dt_val = datetime(*dt_tuple)
                            iso = _date_to_iso(dt_val)
                            if iso and find_text in iso:
                                new_dt = _iso_to_datetime(iso.replace(find_text, replace_text))
                                if new_dt:
                                    value = new_dt
                                    count += 1
                        except Exception:
                            pass
                    if value != '' and value is not None:
                        ws.write(row, col, value)
        if count:
            wb.save(filepath)
        _dbg_write(f"xlwt: count={count}")
        return count
    except ImportError:
        _dbg_write("xlwt FAIL: ImportError")
    except Exception as e:
        _dbg_write(f"xlwt FAIL: {e}")

    _dbg_write("ALL FAIL: return 0")
    return 0


def _replace_single(filepath, ext, find_text, replace_text):
    """对单个文件执行一对替换。返回替换次数，-1 表示缺依赖，异常向上抛出。"""
    if ext == ".docx":
        return replace_in_docx(filepath, find_text, replace_text)
    elif ext == ".doc":
        return replace_in_doc(filepath, find_text, replace_text)
    elif ext == ".xlsx":
        return replace_in_xlsx(filepath, find_text, replace_text)
    elif ext == ".xls":
        return replace_in_xls(filepath, find_text, replace_text)
    elif ext == ".pptx":
        return replace_in_pptx(filepath, find_text, replace_text)
    elif ext in TEXT_EXTENSIONS:
        return replace_in_text_file(filepath, find_text, replace_text)
    else:
        return replace_in_text_file(filepath, find_text, replace_text)


def process_folder(folder, pairs, log_callback=None):
    """
    遍历文件夹，对每个文件依次执行多对替换（按 pairs 顺序）。
    pairs: [(find_text, replace_text), ...]
    返回 {"文件路径": 替换次数, ...}
    """
    results = {}
    total_replaced = 0

    for root, dirs, files in os.walk(folder):
        for filename in files:
            filepath = os.path.join(root, filename)
            ext = Path(filename).suffix.lower()

            if ext in SKIP_EXTENSIONS:
                continue

            if log_callback:
                log_callback(f"处理: {filepath}")

            file_count = 0
            file_error = None

            for find_text, replace_text in pairs:
                if find_text == replace_text or not find_text:
                    continue

                if log_callback:
                    log_callback(f"  [{find_text}  →  {replace_text}]")

                try:
                    cnt = _replace_single(filepath, ext, find_text, replace_text)
                    if cnt is None:
                        cnt = 0
                    if cnt == -1:
                        file_error = "缺少依赖库"
                    elif cnt > 0:
                        file_count += cnt
                except Exception as e:
                    file_error = f"错误: {str(e)[:80]}"
                    break

            if file_error:
                results[filepath] = file_error
            elif file_count > 0:
                results[filepath] = file_count
                total_replaced += file_count

    return results, total_replaced


# ── GUI ─────────────────────────────────────────────────────────────

class PairListFrame(ttk.Frame):
    """多对替换输入的行列表 UI 组件。"""

    def __init__(self, parent):
        super().__init__(parent)
        self._rows = []  # [(find_var, replace_var)]

        header = ttk.Frame(self)
        header.pack(fill=tk.X, pady=(0, 4))
        ttk.Label(header, text="查找文字", width=28).pack(side=tk.LEFT, padx=(0, 4))
        ttk.Label(header, text="", width=2).pack(side=tk.LEFT)
        ttk.Label(header, text="替换为", width=28).pack(side=tk.LEFT, padx=(4, 0))

        self._inner = None
        self.add_row()
        self._rebuild()

    def add_row(self, find_text="", replace_text=""):
        self._rows.append((tk.StringVar(value=find_text), tk.StringVar(value=replace_text)))
        self._rebuild()

    def remove_row(self, index):
        if len(self._rows) <= 1:
            return
        self._rows.pop(index)
        self._rebuild()

    def move_up(self, index):
        if index <= 0:
            return
        self._rows[index], self._rows[index - 1] = self._rows[index - 1], self._rows[index]
        self._rebuild()

    def move_down(self, index):
        if index >= len(self._rows) - 1:
            return
        self._rows[index], self._rows[index + 1] = self._rows[index + 1], self._rows[index]
        self._rebuild()

    def get_pairs(self):
        """返回 [(find, replace), ...]，跳过空查找。"""
        pairs = []
        for fv, rv in self._rows:
            f = fv.get().strip()
            if f:
                pairs.append((f, rv.get()))
        return pairs

    def _rebuild(self):
        if self._inner is not None:
            self._inner.destroy()
        self._inner = ttk.Frame(self)
        self._inner.pack(fill=tk.X)

        for i, (fv, rv) in enumerate(self._rows):
            row_frame = ttk.Frame(self._inner)
            row_frame.pack(fill=tk.X, pady=1)

            ttk.Entry(row_frame, textvariable=fv, width=30).pack(side=tk.LEFT, padx=(0, 4))
            ttk.Label(row_frame, text="→").pack(side=tk.LEFT)
            ttk.Entry(row_frame, textvariable=rv, width=30).pack(side=tk.LEFT, padx=(4, 6))

            if len(self._rows) > 1:
                ttk.Button(row_frame, text="↑", width=2,
                           command=lambda idx=i: self.move_up(idx)).pack(side=tk.LEFT, padx=1)
                ttk.Button(row_frame, text="↓", width=2,
                           command=lambda idx=i: self.move_down(idx)).pack(side=tk.LEFT, padx=1)
                ttk.Button(row_frame, text="✕", width=2,
                           command=lambda idx=i: self.remove_row(idx)).pack(side=tk.LEFT, padx=1)

        ttk.Button(self._inner, text="+ 添加替换对", command=self.add_row).pack(
            fill=tk.X, pady=(6, 0))


class App:
    def __init__(self, root):
        self.root = root
        self.root.title("批量文字替换工具")
        self.root.geometry("720x650")
        self.root.resizable(True, True)

        # 顶部区域 - 文件夹选择
        top = ttk.Frame(root, padding=12)
        top.pack(fill=tk.X)

        ttk.Label(top, text="源文件夹:").pack(side=tk.LEFT)
        self.folder_var = tk.StringVar()
        ttk.Entry(top, textvariable=self.folder_var, width=55).pack(side=tk.LEFT, padx=6)
        ttk.Button(top, text="选择...", command=self.choose_folder).pack(side=tk.LEFT)

        # 多对替换输入
        mid = ttk.Frame(root, padding=12)
        mid.pack(fill=tk.X)

        self.pair_list = PairListFrame(mid)
        self.pair_list.pack(fill=tk.X)

        self.run_btn = ttk.Button(mid, text="开始替换", command=self.start)
        self.run_btn.pack(pady=(12, 0))

        # 进度条
        self.progress = ttk.Progressbar(root, mode="indeterminate")

        # 结果展示
        result_frame = ttk.Frame(root, padding=12)
        result_frame.pack(fill=tk.BOTH, expand=True)

        ttk.Label(result_frame, text="替换报告:").pack(anchor=tk.W)
        self.result_text = tk.Text(result_frame, wrap=tk.WORD, font=("Consolas", 10))
        scroll = ttk.Scrollbar(result_frame, command=self.result_text.yview)
        self.result_text.configure(yscrollcommand=scroll.set)
        self.result_text.pack(fill=tk.BOTH, expand=True, side=tk.LEFT)
        scroll.pack(fill=tk.Y, side=tk.RIGHT)

        # 底部
        ttk.Label(root, text="提示: 操作会先复制源文件夹，在副本中替换，原文件不受影响。",
                  padding=6).pack()

    def choose_folder(self):
        path = filedialog.askdirectory(title="选择要处理的文件夹")
        if path:
            self.folder_var.set(path)

    def log(self, message):
        self.result_text.insert(tk.END, message + "\n")
        self.result_text.see(tk.END)
        self.root.update_idletasks()

    def start(self):
        src = self.folder_var.get().strip()
        pairs = self.pair_list.get_pairs()

        if not src or not os.path.isdir(src):
            messagebox.showerror("错误", "请先选择一个有效的文件夹。")
            return
        if not pairs:
            messagebox.showerror("错误", "请至少添加一对待替换的文字。")
            return

        self.run_btn.configure(state="disabled")
        self.result_text.delete("1.0", tk.END)
        self.progress.pack(fill=tk.X, before=self.result_text.master, pady=(0, 4))
        self.progress.start()

        threading.Thread(target=self._run, args=(src, pairs),
                         daemon=True).start()

    def _run(self, src, pairs):
        # COM 初始化：_run 运行在后台线程，必须初始化 COM apartment
        try:
            import pythoncom
            pythoncom.CoInitialize()
        except ImportError:
            pass

        try:
            parent = os.path.dirname(src.rstrip(os.sep))
            src_name = os.path.basename(src.rstrip(os.sep))
            dst = os.path.join(parent, f"{src_name}_replaced")

            if os.path.exists(dst):
                i = 1
                while os.path.exists(f"{dst}_{i}"):
                    i += 1
                dst = f"{dst}_{i}"

            self.log(f"源文件夹: {src}")
            self.log(f"目标文件夹: {dst}")
            self.log(f"共 {len(pairs)} 对替换:")
            for i, (f, r) in enumerate(pairs, 1):
                self.log(f"  {i}. \"{f}\"  →  \"{r}\"")
            self.log("=" * 60)

            self.log("\n正在复制文件夹...")
            copy_folder(src, dst)
            self.log("复制完成。\n")

            self.log("正在扫描并替换...\n")
            results, total = process_folder(dst, pairs, log_callback=self.log)

            # 4. 汇总报告
            self.log("\n" + "=" * 60)
            self.log("替换完成!\n")

            modified = {k: v for k, v in results.items() if isinstance(v, int) and v > 0}
            errors = {k: v for k, v in results.items() if isinstance(v, str)}

            if modified:
                self.log(f"修改文件数: {len(modified)} 个，共替换 {total} 处\n")
                for path, cnt in sorted(modified.items()):
                    rel = os.path.relpath(path, dst)
                    self.log(f"  {rel}  →  {cnt} 处")
            else:
                self.log("没有找到匹配的文字。")

            if errors:
                self.log(f"\n跳过/出错文件: {len(errors)} 个")
                for path, reason in sorted(errors.items()):
                    rel = os.path.relpath(path, dst)
                    self.log(f"  {rel}  →  {reason}")

        except Exception as e:
            self.log(f"\n发生错误: {e}")
        finally:
            _cleanup_com()
            try:
                import pythoncom
                pythoncom.CoUninitialize()
            except ImportError:
                pass
            self.root.after(0, self._finish)

    def _finish(self):
        self.progress.stop()
        self.progress.pack_forget()
        self.run_btn.configure(state="normal")


if __name__ == "__main__":
    root = tk.Tk()
    app = App(root)
    root.mainloop()
